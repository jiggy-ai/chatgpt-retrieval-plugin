from loguru import logger
import os
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx
import json
from models.models import Document, DocumentMetadata, Source
from services.extract_metadata import extract_metadata_from_document, csv_has_header
import subprocess       
import random
import string
from io import BytesIO 
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
import re






excel_mimetypes = ["application/vnd.ms-excel",                                           # Excel 97-2003 Workbook (.xls)
                   "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # Excel Workbook (.xlsx)
                   "application/vnd.ms-excel.sheet.macroEnabled.12",                     # Excel Macro-Enabled Workbook (.xlsm)
                   "application/vnd.ms-excel.sheet.binary.macroEnabled.12",              # Excel Binary Workbook (.xlsb) 
                   "application/vnd.ms-excel.template.macroEnabled.12",                  # Excel Template (.xlt)
                   "application/vnd.ms-excel.template.macroEnabled.12",                  # Excel Macro-Enabled Template (.xltm)
                   "application/vnd.ms-excel.addin.macroEnabled.12"]                     # Excel Add-In (.xlam)


async def get_document_from_file(file: UploadFile, id: str = None, metadata: Optional[DocumentMetadata] = None) -> Document:
    extracted_text, mimetype = await extract_text_from_form_file(file)

    if mimetype not in excel_mimetypes + ['text/csv']:
        extracted_metadata = extract_metadata_from_document(extracted_text, file.filename)
        logger.info(f"Extracted metadata: {extracted_metadata}")
    else:
        # tabular datafiles tend not to have visible metadata
        extracted_metadata = {}        
        logger.info("No metadata extracted for excel or csv files")

    if metadata:
        # update metadata keys that are None using the extracted metadata
        metadata_dict = metadata.dict()
        for k, v in extracted_metadata.items():
            if v and k in metadata_dict and metadata_dict[k] is None:
                metadata.__setattr__(k, v)
    else: 
        metadata = DocumentMetadata(**extracted_metadata)                                     

    # set source and source_id if not already set                                    
    if metadata.source is None:
        metadata.source = Source.file
    if metadata.source_id is None:
        metadata.source_id = file.filename
        
    logger.info(metadata)
    doc = Document(text=extracted_text, metadata=metadata, mimetype=mimetype)
    if id:
        doc.id = id
    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if not mimetype:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        logger.info(f"No mimetype found for {filepath}, infer based on file extension")
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        elif filepath.endswith(".js"):
            mimetype = "application/javascript"
        elif filepath.endswith(".ts"):
            mimetype = "application/x-typescript"
        elif filepath.endswith(".tsx"):
            mimetype = "application/x-typescript-jsx"
        else:
            logger.info(f"Unable to infer mimetype for {filepath}")
            raise ValueError("Unsupported file type")

    # Open the file in binary mode
    with open(filepath, "rb") as file:
        return extract_text_from_file(file, mimetype)




# dictionary of excel mimetypes and their corresponding extensions

excel_mimetypes_to_extensions = {
    "application/vnd.ms-excel": "xls",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel.sheet.macroEnabled.12": "xlsm",
    "application/vnd.ms-excel.sheet.binary.macroEnabled.12": "xlsb",
    "application/vnd.ms-excel.template.macroEnabled.12": "xlt",
    "application/vnd.ms-excel.template.macroEnabled.12": "xltm",
    "application/vnd.ms-excel.addin.macroEnabled.12": "xlam",
}

encoding_options = (("utf-8", "strict"), ("windows-1252", "strict"), ("iso-8859-1", "strict"), ("utf-8", "replace"))

def dynamic_decode(data):
    for encoding, errors in encoding_options:
        try:
            return data.decode(encoding, errors=errors)
        except UnicodeDecodeError:
            pass
    raise ValueError("Unable to decode data, unknown character encoding")


def csv_fp_to_text(file) -> str:
    extracted_text = ""    
    decoded_buffer = [dynamic_decode(line) for line in file]
    if csv_has_header(decoded_buffer[:20]): 
        logger.info("reading csv with header")
        reader = csv.DictReader(decoded_buffer)
        for row in reader:
            extracted_text += json.dumps(row) + "\n"
    else:
        logger.info("reading csv without header")
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"    
    return extracted_text


text_mimetypes = ["text/plain", 
                  "text/markdown", 
                  "text/x-python",
                  "application/x-typescript",
                  "application/x-typescript-jsx",
                  "application/javascript"]


def pdf_text(file):
    """
    extract text from the pdf_bytes using pdfminer.six and return it as a single string
    """
    text = ""
    for page_layout in extract_pages(file):
        for element in page_layout:
            if isinstance(element, LTTextContainer):
                for text_line in element:
                    line = text_line.get_text().rstrip() + "\n"
                    text += line
    return re.sub(r'\n{3,}', '\n\n', text)


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    
    if mimetype == "application/pdf":
        try:
            extracted_text = pdf_text(file)
            if not extracted_text or len(extracted_text.strip()) < 10:
                raise Exception("pdfminer.six failed to extract useful text from pdf")
            logger.info("Extracted text from pdf using pdfminer.six")
        except Exception as e:
            logger.error(f'Failed to extract text from pdf using pdfminer.six ({e})')
            # Extract text from pdf using PyPDF2
            reader = PdfReader(file)
            extracted_text = " ".join([page.extract_text() for page in reader.pages])
            logger.info("Extracted text from pdf using PyPDF2")
            
    elif mimetype in text_mimetypes:
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
        
    elif (mimetype == "application/msword"):
        input_file = "/tmp/tmp.doc"
        open(input_file, 'wb').write(file.read())
        output_folder = "/tmp/"
        output_format = "docx:\"Office Open XML Text\""
        command = f"libreoffice --headless --convert-to {output_format} --outdir {output_folder} {input_file}"
        try:
            subprocess.run(command, shell=True, check=True)
        except:
            raise ValueError("Unable to convert doc to docx.")      
        extracted_text = docx2txt.process(f"{output_folder}/tmp.docx")  
        os.unlink(f"{output_folder}/tmp.docx")
        os.unlink(input_file)
        
    elif (mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
        
    elif mimetype in excel_mimetypes_to_extensions:
        input_file = f"/tmp/tmp.{excel_mimetypes_to_extensions[mimetype]}"
        open(input_file, 'wb').write(file.read())
        output_folder = "/tmp/"
        command = f"libreoffice --headless --convert-to csv --outdir {output_folder} {input_file}"
        try:
            subprocess.run(command, shell=True, check=True)
        except:
            raise ValueError("Unable to convert excel to csv.")      
        with open(f"{output_folder}/tmp.csv", 'rb') as csv_file:
            extracted_text = csv_fp_to_text(csv_file)
        os.unlink(f"{output_folder}/tmp.csv")
        os.unlink(input_file)
        
    elif mimetype == "text/csv":
        extracted_text = csv_fp_to_text(file)
        
    elif (mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation"):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text



def infer_mimetype(filename):
    mimetype, _ = mimetypes.guess_type(filename)
    
    if not mimetype:
        logger.info(f"No mimetype found for {filename}, infer based on file extension")
        if filename.endswith(".md"):
            mimetype = "text/markdown"
        elif filename.endswith(".js"):
            mimetype = "application/javascript"
        elif filename.endswith(".ts"):
            mimetype = "application/x-typescript"
        elif filename.endswith(".tsx"):
            mimetype = "application/x-typescript-jsx"
        else:
            logger.info(f"Unable to infer mimetype for {filename}")
            raise ValueError("Unknown or unsupported file type. Try specifying the mimetype.")
    return mimetype


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    if not mimetype:
        mimetype = infer_mimetype(file.filename)
    logger.info(f"mimetype: {mimetype}")

    file_stream = await file.read()

    temp_file_path = "/tmp/" + ''.join(random.choices(string.ascii_letters, k=8))

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        with open(temp_file_path, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype)        
        os.remove(temp_file_path)
    except ValueError as e:
        raise
    except Exception as e:
        logger.exception(f"Error extracting text from file: {e}")
        raise ValueError(f"Error extracting text from file")
                
    return extracted_text, mimetype
